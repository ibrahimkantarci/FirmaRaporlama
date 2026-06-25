# api/generate-deck.py — Vercel Python serverless fonksiyonu (Next.js ile birlikte).
# Düzenleme ekranından gelen JSON'ı template.pptx üstüne basıp .pptx döndürür.
# Saf python-pptx (skill scriptlerine bağımsız). Endpoint: POST /api/generate-deck
from http.server import BaseHTTPRequestHandler
import copy, io, json, os
from pptx import Presentation
from pptx.util import Inches
from pptx.oxml.ns import qn
from pptx.opc.constants import RELATIONSHIP_TYPE as RT

A = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
P = "{http://schemas.openxmlformats.org/presentationml/2006/main}"
R = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}"

TEMPLATE = os.path.join(os.path.dirname(__file__), "template.pptx")

# --- Template'e özgü anahtarlar (template.pptx incelemesinden) ---
GROUP_NAME = "Google Shape;136;p29"   # slayt-3 chevron grubu
BANNER_L = "Google Shape;132;p29"
BANNER_R = "Google Shape;133;p29"
COVER_KEY = "Portaxe"
ORIG = {
    "win": "Winner 6x", "title": "Portaxe-Balo ve Davet Salonları", "dash": "-",
    "sayfa": "19.263", "cift": "1.325", "donus": " 8,5 Saat", "profil": "98",
}
BAR_MAP = {
    "Google Shape;156;p30": ("Google Shape;157;p30", "sayfa"),
    "Google Shape;168;p30": ("Google Shape;169;p30", "teklif"),
    "Google Shape;180;p30": ("Google Shape;181;p30", "donus"),
    "Google Shape;192;p30": ("Google Shape;193;p30", "profil"),
}
TRACK = Inches(2.72)


def _num(x):
    s = str(x).strip().replace(",", "")
    if s in ("", "nan", "None"):
        return 0.0
    try:
        return float(s)
    except Exception:
        return 0.0

def tr_int(v):
    return f"{int(round(v)):,}".replace(",", ".")

def tr_dec1(v):
    return f"{v:.1f}".replace(".", ",")

# Dönüş süresi: saat cinsinden değer -> "X saat Y dakika" (sıfır parça atılır).
def hm(v):
    total = int(round(_num(v) * 60))
    h, m = divmod(total, 60)
    if h and m:
        return f"{h} saat {m} dakika"
    if h:
        return f"{h} saat"
    return f"{m} dakika"

def pct(cur, prev):
    if prev == 0:
        return "▲ 0.0%"
    d = (cur - prev) / prev * 100
    arrow = "▲" if cur >= prev else "▼"
    return f"{arrow} {abs(d):.1f}%"


# --- slayt klonlama (saf python-pptx) ---
def _remap_rels(src_part, dst_part):
    rid_map = {}
    for rId, rel in src_part.rels.items():
        if rel.reltype == RT.SLIDE_LAYOUT:
            continue
        if rel.is_external:
            new_rId = dst_part.rels.get_or_add_ext_rel(rel.reltype, rel.target_ref)
        else:
            new_rId = dst_part.relate_to(rel.target_part, rel.reltype)
        rid_map[rId] = new_rId
    return rid_map

def _apply_rid_map(spTree, rid_map):
    for el in spTree.iter():
        for attr in (f"{R}embed", f"{R}link", f"{R}id"):
            v = el.get(attr)
            if v is not None and v in rid_map:
                el.set(attr, rid_map[v])

def clone_slide(prs, src_slide):
    new_slide = prs.slides.add_slide(src_slide.slide_layout)
    for shp in list(new_slide.shapes):
        shp._element.getparent().remove(shp._element)
    src_spTree = src_slide.shapes._spTree
    new_spTree = new_slide.shapes._spTree
    for child in list(src_spTree):
        if child.tag in (f"{P}nvGrpSpPr", f"{P}grpSpPr"):
            continue
        new_spTree.append(copy.deepcopy(child))
    rid_map = _remap_rels(src_slide.part, new_slide.part)
    _apply_rid_map(new_spTree, rid_map)
    return new_slide

def reorder_slides(prs, desired):
    sldIdLst = prs.slides._sldIdLst
    id_to_el = {el.get(qn("r:id")): el for el in list(sldIdLst)}
    part_to_rid = {}
    for rId, rel in prs.part.rels.items():
        if rel.reltype == RT.SLIDE and not rel.is_external:
            part_to_rid[rel.target_part] = rId
    for el in list(sldIdLst):
        sldIdLst.remove(el)
    for s in desired:
        sldIdLst.append(id_to_el[part_to_rid[s.part]])


# --- metin doldurma ---
def _set_texts(el, mapping):
    done = {k: False for k in mapping}
    for t in el.iter(f"{A}t"):
        tx = t.text or ""
        if tx in mapping and not done[tx]:
            t.text = mapping[tx]
            done[tx] = True

def _find_shape(slide, name):
    for sh in slide.shapes:
        if sh.name == name:
            return sh
    return None

def _grp_box(grp_el):
    xf = grp_el.find(f"{P}grpSpPr/{A}xfrm")
    return xf.find(f"{A}off"), xf.find(f"{A}ext")

def _fill_one(grp_el, c):
    _set_texts(grp_el, {
        ORIG["win"]: str(c.get("urun", "")),
        ORIG["title"]: f'{c.get("rci","")} - {c.get("kategori","")}',
        ORIG["dash"]: "",
        ORIG["sayfa"]: tr_int(_num(c.get("sayfa"))),
        ORIG["cift"]: tr_int(_num(c.get("teklif"))),
        ORIG["donus"]: f' {hm(c.get("donus"))}',
        ORIG["profil"]: tr_int(_num(c.get("profil"))),
    })

def _fill_banner(slide, venue):
    # "Her zaman göster, boş alanlar": madde yoksa şablonun ön-yazılı metni TEMİZLENİR.
    bullets = list(venue.get("bannerBullets") or [])
    left, right = bullets[:3], bullets[3:6]
    for name, items in ((BANNER_L, left), (BANNER_R, right)):
        sh = _find_shape(slide, name)
        if sh is None or not sh.has_text_frame:
            continue
        for i, p in enumerate(sh.text_frame.paragraphs):
            new = str(items[i]) if i < len(items) else ""
            runs = p.runs
            if runs:
                runs[0].text = new
                for r in runs[1:]:
                    r.text = ""
            elif new:
                p.text = new

def _fill_venue(slide, venue):
    rci = venue.get("rci", "")
    cats = [{**c, "rci": rci} for c in venue.get("categories", [])]
    grp = _find_shape(slide, GROUP_NAME)
    if grp is None:
        return
    gel = grp._element
    K = len(cats)
    if K <= 1:
        if K == 1:
            _fill_one(gel, cats[0])
        _fill_banner(slide, venue)
        return
    off, ext = _grp_box(gel)
    orig_cy = int(ext.get("cy"))
    y_start = Inches(1.30)
    banner_top = Inches(4.12)
    avail = banner_top - y_start
    row_h = min(orig_cy, int(avail / K))
    snap = copy.deepcopy(gel)
    groups = [gel]
    prev = gel
    for _ in range(1, K):
        gg = copy.deepcopy(snap)
        prev.addnext(gg)
        groups.append(gg)
        prev = gg
    for i, (gg, c) in enumerate(zip(groups, cats)):
        _fill_one(gg, c)
        o, e = _grp_box(gg)
        o.set("y", str(int(y_start + (avail / K) * i)))
        e.set("cy", str(int(row_h)))
    _fill_banner(slide, venue)


# --- totals ---
def _aggregate(venues):
    cats = [c for v in venues for c in v.get("categories", [])]
    def s(key):
        return sum(_num(c.get(key)) for c in cats)
    def a(key):
        xs = [_num(c.get(key)) for c in cats if _num(c.get(key)) > 0]
        return (sum(xs) / len(xs)) if xs else 0.0
    return {
        "sayfa": (s("sayfa"), s("sayfaGy")),
        "teklif": (s("teklif"), s("teklifGy")),
        "donus": (a("donus"), a("donusGy")),
        "profil": (a("profil"), a("profilGy")),
    }

def _fill_totals(slide, agg):
    sayfa, teklif, donus, profil = agg["sayfa"], agg["teklif"], agg["donus"], agg["profil"]
    t_repl = {
        "▼ 2.0%": pct(*sayfa), "19.658": tr_int(sayfa[1]), "19.263": tr_int(sayfa[0]),
        "▲ 11.6%": pct(*teklif), "1.187": tr_int(teklif[1]), "1.325": tr_int(teklif[0]),
        "▲ 30.8%": pct(*donus), "6,5 saat": hm(donus[1]), "8,5 saat": hm(donus[0]),
        "▲ 22.5%": pct(*profil), "80": tr_int(profil[1]), "98": tr_int(profil[0]),
    }
    done = {k: False for k in t_repl}
    for t in slide.shapes._spTree.iter(f"{A}t"):
        tx = t.text or ""
        if tx in t_repl and not done[tx]:
            t.text = t_repl[tx]
            done[tx] = True
    shp = {sh.name: sh for sh in slide.shapes}
    for gname, (mname, key) in BAR_MAP.items():
        cur, prev = agg[key]
        mx = max(cur, prev) or 1
        if gname in shp:
            shp[gname].width = int(TRACK * (prev / mx))
        if mname in shp:
            shp[mname].width = int(TRACK * (cur / mx))


def build_deck(data, template_path=TEMPLATE):
    prs = Presentation(template_path)
    slides = list(prs.slides)
    cover, intro, venue_tpl, totals, advice, closing = slides[:6]
    venues = data.get("venues", [])
    if not venues:
        raise ValueError("venue yok")
    cover_name = data.get("coverName") or "Firma"
    for t in cover.shapes._spTree.iter(f"{A}t"):
        if (t.text or "") == COVER_KEY:
            t.text = cover_name
            break
    venue_slides = [venue_tpl]
    for _ in range(1, len(venues)):
        venue_slides.append(clone_slide(prs, venue_tpl))
    for slide, venue in zip(venue_slides, venues):
        _fill_venue(slide, venue)
    _fill_totals(totals, _aggregate(venues))
    reorder_slides(prs, [cover, intro] + venue_slides + [totals, advice, closing])
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


class handler(BaseHTTPRequestHandler):
    def _err(self, code, msg):
        self.send_response(code)
        self.send_header("Content-Type", "application/json; charset=utf-8")
        self.end_headers()
        self.wfile.write(json.dumps({"error": msg}, ensure_ascii=False).encode("utf-8"))

    def do_POST(self):
        try:
            length = int(self.headers.get("content-length", 0) or 0)
            body = self.rfile.read(length) if length else b"{}"
            data = json.loads(body.decode("utf-8"))
        except Exception as e:
            return self._err(400, f"Geçersiz JSON: {e}")
        try:
            pptx = build_deck(data)
        except Exception as e:
            return self._err(500, f"Üretim hatası: {e}")
        self.send_response(200)
        self.send_header(
            "Content-Type",
            "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        )
        self.send_header("Content-Disposition", 'attachment; filename="deck.pptx"')
        self.send_header("Content-Length", str(len(pptx)))
        self.end_headers()
        self.wfile.write(pptx)
