#!/usr/bin/env python3
# Armonia bloğundan deck üretir.
import copy, subprocess, sys, re
from pathlib import Path
import pandas as pd
from lxml import etree
from pptx import Presentation
from pptx.util import Inches, Emu

A = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
P = "{http://schemas.openxmlformats.org/presentationml/2006/main}"
R = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}"
PPTX = "/home/claude/template.pptx"
XLSX = "/mnt/user-data/uploads/Firma_Raporlama.xlsx"
UNP = Path("/home/claude/gen_unpacked")
SKILL = "/mnt/skills/public/pptx/scripts"

# ---------- 1) Excel: Armonia bloğunu ayıkla ----------
df = pd.read_excel(XLSX, header=None)
n = len(df)
def cell(i, j):
    v = df.iloc[i, j]
    return "" if (pd.isna(v)) else str(v)

# blok başlangıçları (col0 "Müşteri:")
starts = [i for i in range(n) if cell(i, 0).startswith("Müşteri:")]
starts.append(n)
block = None
for bi in range(len(starts) - 1):
    s = starts[bi]; e = starts[bi + 1]
    meta = {cell(s, k).split(":")[0].strip(): cell(s, k).split(":", 1)[1].strip()
            for k in range(df.shape[1]) if ":" in cell(s, k)}
    header = [cell(s + 1, k) for k in range(df.shape[1])]
    rows = [df.iloc[r] for r in range(s + 2, e)]
    if "Armonia" in cell(s + 2, 4):  # Müşteri Adı sütununda Armonia
        block = (meta, header, rows); break
if not block:
    block = (meta, header, rows)
meta, header, rows = block

def col(name):
    return header.index(name)
C = {
    "rci": col("RÇİ Adı"), "kat": col("Kategori Adı"), "must": col("Müşteri Adı"), "urun": col("Ürün Adı"),
    "sayfa": col("Sayfa Ziyareti"), "teklif": col("Teklif"),
    "donus": col("Ortalama Dönüş Süresi (Saat)"), "profil": col("Profil Puanı"),
    "sayfa_gy": col("Sayfa Ziyareti (GY)"), "teklif_gy": col("Teklif (GY)"),
    "donus_gy": col("Ortalama Dönüş Süresi (Saat) (GY)"), "profil_gy": col("Profil Puanı (GY)"),
}

def num(x):
    s = str(x).strip().replace(",", "")  # virgül = binlik
    if s in ("", "nan", "None"):
        return 0.0
    try:
        return float(s)
    except:
        return 0.0

def tr_int(v):
    return f"{int(round(v)):,}".replace(",", ".")
def tr_dec1(v):
    return f"{v:.1f}".replace(".", ",")

firm = (meta.get("Müşteri") and "") or ""
firm_name = cell(0, 4) if False else ""  # placeholder
# Kapak adı: Müşteri Adı'nın ilk kelimesi
must_full = str(rows[0][C["must"]])
cover_name = must_full.split()[0] if must_full else "Firma"
this_date = meta.get("Bu yıl", "")

# venue gruplama (RÇİ Adı), sırayı koru
from collections import OrderedDict
venues = OrderedDict()
for r in rows:
    v = str(r[C["rci"]])
    venues.setdefault(v, []).append(r)
venue_list = list(venues.keys())
print(f"Müşteri: {cover_name} | venue: {len(venue_list)} | bu dönem: {this_date}")

# ---------- 2) Template'i aç, slayt 3'ü venue başına çoğalt ----------
import shutil
if UNP.exists(): shutil.rmtree(UNP)
subprocess.run([sys.executable, f"{SKILL}/office/unpack.py", PPTX, str(UNP)+"/"], check=True,
               capture_output=True)

need_copies = len(venue_list) - 1  # slide3 zaten 1 venue
new_sldids = []
for _ in range(need_copies):
    out = subprocess.run([sys.executable, f"{SKILL}/add_slide.py", str(UNP), "slide3.xml"],
                         check=True, capture_output=True, text=True)
    m = re.search(r'<p:sldId\s[^>]*/>', out.stdout)
    new_sldids.append(m.group(0))

# presentation.xml: sldIdLst'i istediğimiz sıraya diz
pres_path = UNP / "ppt" / "presentation.xml"
ns = {"p": P[1:-1], "r": R[1:-1]}
tree = etree.parse(str(pres_path))
root = tree.getroot()
sldIdLst = root.find(f"{P}sldIdLst")
orig = list(sldIdLst)  # 6 orijinal: s1,s2,s3,s4,s5,s6
# rel map: rId -> slide dosyası
rels = etree.parse(str(UNP/"ppt"/"_rels"/"presentation.xml.rels")).getroot()
rid2file = {rel.get("Id"): rel.get("Target") for rel in rels}
def file_of(sldid_el):
    return rid2file.get(sldid_el.get(f"{R}id"), "")
orig_by_file = {file_of(s).split("/")[-1]: s for s in orig}
# yeni sldId elemanlarını parse et
_wrap = lambda x: etree.fromstring(f'<root xmlns:p="{P[1:-1]}" xmlns:r="{R[1:-1]}">{x}</root>')[0]
new_els = [_wrap(x) for x in new_sldids]
# istenen sıra: s1, s2, s3(venue1), yeni1..N(venue2..), s4, s5, s6
order = [orig_by_file["slide1.xml"], orig_by_file["slide2.xml"], orig_by_file["slide3.xml"]]
order += new_els
order += [orig_by_file["slide4.xml"], orig_by_file["slide5.xml"], orig_by_file["slide6.xml"]]
# id attribute'larını benzersiz yap
for i, el in enumerate(order):
    el.set("id", str(256 + i))
for s in list(sldIdLst):
    sldIdLst.remove(s)
for el in order:
    sldIdLst.append(el)
tree.write(str(pres_path), xml_declaration=True, encoding="UTF-8", standalone=True)

subprocess.run([sys.executable, f"{SKILL}/clean.py", str(UNP)], capture_output=True)
subprocess.run([sys.executable, f"{SKILL}/office/pack.py", str(UNP), "/home/claude/deck_raw.pptx",
                "--original", PPTX], check=True, capture_output=True)
print("deck_raw.pptx hazır, slayt sayısı düzenlendi")

# ---------- 3) python-pptx ile doldur ----------
prs = Presentation("/home/claude/deck_raw.pptx")
slides = list(prs.slides)
# index: 0 kapak, 1 intro, 2..(2+V-1) venue, sonra totals, advice, closing
V = len(venue_list)
cover = slides[0]
venue_slides = slides[2:2+V]
totals = slides[2+V]

def set_texts(el, mapping):
    done = {k: False for k in mapping}
    for t in el.iter(f"{A}t"):
        tx = t.text or ""
        if tx in mapping and not done[tx]:
            t.text = mapping[tx]; done[tx] = True

# kapak adı
for t in cover.shapes._spTree.iter(f"{A}t"):
    if (t.text or "") == "Portaxe":
        t.text = cover_name

# --- slayt 3 şablon metin anahtarları (orijinal Portaxe satırı) ---
ORIG = {"title": "Portaxe-Balo ve Davet Salonları", "dash": "-", "win": "Winner 6x",
        "sayfa": "19.263", "cift": "1.325", "donus": " 8,5 Saat", "profil": "98"}
BULLET_NAMES = {"Google Shape;131;p29","Google Shape;132;p29","Google Shape;133;p29",
                "Google Shape;135;p29","Google Shape;134;p29"}

def find_group(slide):
    for sh in slide.shapes:
        if sh.name == "Google Shape;136;p29":
            return sh._element
    return None
def grp_top(el, emu):
    off = el.find(f"{P}grpSpPr/{A}xfrm/{A}off"); off.set("y", str(int(emu)))

def grp_box(el):
    xf = el.find(f"{P}grpSpPr/{A}xfrm")
    return xf.find(f"{A}off"), xf.find(f"{A}ext")

def fill_one(g, r):
    set_texts(g, {
        ORIG["win"]: str(r[C["urun"]]),                       # Ürün Adı (örn. Winner 6X)
        ORIG["title"]: f'{r[C["rci"]]} - {r[C["kat"]]}',      # RÇİ Adı - Kategori
        ORIG["dash"]: "",
        ORIG["sayfa"]: tr_int(num(r[C["sayfa"]])),
        ORIG["cift"]: tr_int(num(r[C["teklif"]])),
        ORIG["donus"]: f" {tr_dec1(num(r[C['donus']]))} Saat",
        ORIG["profil"]: tr_int(num(r[C["profil"]])),
    })

def fill_venue(slide, venue, cats):
    K = len(cats)
    grp = find_group(slide)
    if K == 1:
        # template gibi: tek satır, banner durur
        fill_one(grp, cats[0])
        return
    # çok kategori: banner KORUNUR; satırlar üstüne sığar (gerekirse ölçekle)
    off, ext = grp_box(grp)
    orig_cy = int(ext.get("cy"))
    y_start = Inches(1.30)
    banner_top = Inches(4.12)
    avail = banner_top - y_start
    row_h = min(orig_cy, int(avail / K))
    snap = copy.deepcopy(grp)
    groups = [grp]; prev = grp
    for i in range(1, K):
        gg = copy.deepcopy(snap); prev.addnext(gg); groups.append(gg); prev = gg
    for i, (gg, r) in enumerate(zip(groups, cats)):
        fill_one(gg, r)
        o, e = grp_box(gg)
        o.set("y", str(int(y_start + (avail / K) * i)))
        e.set("cy", str(int(row_h)))   # yüksekliği sığacak şekilde ölçekle (çocuklar küçülür)

for slide, venue in zip(venue_slides, venue_list):
    fill_venue(slide, venue, venues[venue])

# --- totals slaytı (slayt 4) ---
def agg(col_this, col_gy, mode):
    vals = [num(r[col_this]) for r in rows]
    valg = [num(r[col_gy]) for r in rows]
    if mode == "sum":
        return sum(vals), sum(valg)
    else:  # avg (0'ları say ama gerçekçi tut)
        vv = [x for x in vals if x > 0]; gg = [x for x in valg if x > 0]
        return (sum(vv)/len(vv) if vv else 0), (sum(gg)/len(gg) if gg else 0)

m_sayfa = agg(C["sayfa"], C["sayfa_gy"], "sum")
m_cift  = agg(C["teklif"], C["teklif_gy"], "sum")
m_donus = agg(C["donus"], C["donus_gy"], "avg")
m_profil= agg(C["profil"], C["profil_gy"], "avg")

def pct(cur, prev):
    if prev == 0: return "▲ 0.0%"
    d = (cur - prev) / prev * 100
    arrow = "▲" if cur >= prev else "▼"
    return f"{arrow} {abs(d):.1f}%"

# slayt 4 metin değişimleri (şablon değerleri -> yeni)
t_repl = {
    "▼ 2.0%": pct(m_sayfa[0], m_sayfa[1]), "19.658": tr_int(m_sayfa[1]), "19.263": tr_int(m_sayfa[0]),
    "▲ 11.6%": pct(m_cift[0], m_cift[1]),  "1.187": tr_int(m_cift[1]),  "1.325": tr_int(m_cift[0]),
    "▲ 30.8%": pct(m_donus[0], m_donus[1]),"6,5 saat": tr_dec1(m_donus[1])+" saat", "8,5 saat": tr_dec1(m_donus[0])+" saat",
    "▲ 22.5%": pct(m_profil[0], m_profil[1]),"80": tr_int(m_profil[1]), "98": tr_int(m_profil[0]),
}
done = {k: False for k in t_repl}
for t in totals.shapes._spTree.iter(f"{A}t"):
    tx = t.text or ""
    if tx in t_repl and not done[tx]:
        t.text = t_repl[tx]; done[tx] = True

# orantılı barlar: value bar genişliği = track * v/max
TRACK = Inches(2.72)
bar_map = {  # (gecen_bar, mevcut_bar, gecen_val, mevcut_val)
    "Google Shape;156;p30": ("Google Shape;157;p30", m_sayfa),
    "Google Shape;168;p30": ("Google Shape;169;p30", m_cift),
    "Google Shape;180;p30": ("Google Shape;181;p30", m_donus),
    "Google Shape;192;p30": ("Google Shape;193;p30", m_profil),
}
shp = {sh.name: sh for sh in totals.shapes}
for gname, (mname, (cur, prev)) in bar_map.items():
    mx = max(cur, prev) or 1
    shp[gname].width = int(TRACK * (prev / mx))
    shp[mname].width = int(TRACK * (cur / mx))

prs.save("/home/claude/deck.pptx")
print("deck.pptx üretildi.")
